"""Build the NeuroMesh submission deck as an editable .pptx.

    pip install python-pptx
    python tools/make_deck_pptx.py

Everything is a native PowerPoint object -- text boxes, tables, autoshapes --
so every word, colour and box can be edited in PowerPoint or Google Slides.
Nothing is a flattened image.

Fonts are Segoe UI and Consolas because they ship with Windows and Office. If
you want the web deck's typography, install Archivo and IBM Plex Mono and change
DISPLAY / MONO below.
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

OUT = pathlib.Path(__file__).resolve().parent.parent / "NeuroMesh_iQOO_Hackathon.pptx"

# ---------------------------------------------------------------- palette ---
GROUND = RGBColor(0x07, 0x0B, 0x14)
PANEL = RGBColor(0x11, 0x18, 0x28)
PANEL2 = RGBColor(0x16, 0x1E, 0x31)
EDGE = RGBColor(0x26, 0x30, 0x4A)
INK = RGBColor(0xE9, 0xED, 0xF7)
HAZE = RGBColor(0x8B, 0x95, 0xB3)
MINT = RGBColor(0x4F, 0xE0, 0xBC)
IRIS = RGBColor(0x8A, 0xA5, 0xFF)
EMBER = RGBColor(0xFF, 0x84, 0x74)
IDLE = RGBColor(0x2A, 0x33, 0x50)
MINT_DIM = RGBColor(0x1B, 0x4A, 0x44)
IRIS_DIM = RGBColor(0x24, 0x30, 0x5C)
EMBER_DIM = RGBColor(0x4A, 0x24, 0x22)

DISPLAY = "Segoe UI Semibold"
BODY = "Segoe UI"
MONO = "Consolas"

# ------------------------------------------------------------ geometry -----
W, H = 13.333, 7.5
L = 0.55                      # left margin
CW = W - 2 * L                # content width


# ------------------------------------------------------------- helpers -----
def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts -> text, size, color, font, bold, space_after, line."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if spec.get("space_before"):
            p.space_before = Pt(spec["space_before"])
        p.space_after = Pt(spec.get("space_after", 0))
        if spec.get("line"):
            p.line_spacing = spec["line"]
        run = p.add_run()
        run.text = spec["text"]
        f = run.font
        f.size = Pt(spec.get("size", 11))
        f.name = spec.get("font", BODY)
        f.bold = spec.get("bold", False)
        f.color.rgb = spec.get("color", INK)
        if spec.get("spacing"):
            run.font._rPr.set("spc", str(int(spec["spacing"] * 100)))
    return box


def box(slide, x, y, w, h, fill=PANEL, line=EDGE, radius=0.045, line_w=0.75):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    shp.text_frame.word_wrap = True
    return shp


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def line(slide, x1, y1, x2, y2, color=HAZE, width=1.0, dashed=False, arrow=True):
    """A connector with a real arrowhead (set through the drawing XML)."""
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        dash = etree.SubElement(ln, qn("a:prstDash"))
        dash.set("val", "dash")
    if arrow:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return conn


def slide_base(prs, index, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GROUND

    textbox(slide, L, 0.34, CW, 0.24, [{
        "text": f"{index:02d} / 06     {section.upper()}",
        "size": 9, "font": MONO, "color": HAZE, "spacing": 1.6,
    }])
    rect(slide, L, 0.63, CW, 0.012, EDGE)
    return slide


def title(slide, text, y=0.82, size=30, w=None, color=INK):
    return textbox(slide, L, y, w or CW, 1.0, [{
        "text": text, "size": size, "font": DISPLAY, "bold": True,
        "color": color, "line": 0.95,
    }])


def card(slide, x, y, w, h, label=None, heading=None, body=None, plain=None,
         bullets=None, stat=None, stat_color=MINT):
    """The repeated panel: small label, heading, body, optional plain-English line."""
    box(slide, x, y, w, h)
    pad = 0.16
    cy = y + pad
    inner = w - 2 * pad

    if label:
        textbox(slide, x + pad, cy, inner, 0.18, [{
            "text": label.upper(), "size": 7.5, "font": MONO, "color": HAZE, "spacing": 1.3,
        }])
        cy += 0.26
    if stat:
        textbox(slide, x + pad, cy, inner, 0.42, [{
            "text": stat, "size": 26, "font": DISPLAY, "bold": True, "color": stat_color,
        }])
        cy += 0.52
    if heading:
        n = max(1, len(heading) // 34 + 1)
        textbox(slide, x + pad, cy, inner, 0.24 * n, [{
            "text": heading, "size": 11.5, "font": DISPLAY, "bold": True, "color": INK,
            "line": 1.05,
        }])
        cy += 0.21 * n + 0.10
    if body:
        n = max(1, len(body) // 52 + 1)
        textbox(slide, x + pad, cy, inner, 0.17 * n, [{
            "text": body, "size": 9, "color": HAZE, "line": 1.22,
        }])
        cy += 0.155 * n + 0.10
    if bullets:
        specs = []
        for b in bullets:
            specs.append({
                "text": "— " + b, "size": 9, "color": HAZE,
                "line": 1.2, "space_after": 4,
            })
        avail = h - (cy - y) - pad - (0.52 if plain else 0)
        textbox(slide, x + pad, cy, inner, avail, specs)
        cy += 0.4
    if plain:
        py = y + h - pad - 0.42
        rect(slide, x + pad, py, 0.022, 0.40, MINT)
        textbox(slide, x + pad + 0.11, py, inner - 0.11, 0.40, [
            {"text": "PLAINLY", "size": 6.5, "font": MONO, "bold": True,
             "color": MINT, "spacing": 1.3, "space_after": 1},
            {"text": plain, "size": 8.5, "color": INK, "line": 1.18},
        ])


def chip(slide, x, y, text, color=HAZE, edge=EDGE, w=None):
    w = w or (0.075 * len(text) + 0.24)
    shp = box(slide, x, y, w, 0.28, fill=None, line=edge, radius=0.5, line_w=0.75)
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(8)
    r.font.name = MONO
    r.font.color.rgb = color
    return x + w + 0.12


def table(slide, x, y, w, h, rows, widths, header_color=HAZE):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False

    for i, frac in enumerate(widths):
        tbl.columns[i].width = Emu(int(Inches(w) * frac))

    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(h / len(rows))
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GROUND if r else PANEL2
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            highlight = isinstance(val, tuple)
            text = val[0] if highlight else val

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            f = run.font
            if r == 0:
                f.size, f.name, f.color.rgb, f.bold = Pt(7.5), MONO, header_color, False
            elif highlight:
                f.size, f.name = Pt(9.5), BODY
                f.bold = True
                f.color.rgb = MINT if c == 0 else INK
            else:
                f.size, f.name = Pt(9.5), BODY
                f.bold = c == 0
                f.color.rgb = INK if c == 0 else HAZE
    return tbl


# =========================================================== the slides =====
def slide1(prs):
    s = slide_base(prs, 1, "Problem & Idea")

    title(s, "Your phone has an AI chip.\nEvery night, it does nothing.", y=0.82, size=34)

    rect(s, L, 2.30, 3.2, 0.022, MINT)

    textbox(s, L, 2.52, 9.4, 1.30, [
        {"text": "Modern phones carry a dedicated AI chip — an NPU — capable of roughly "
                 "40 trillion operations per second. From the moment you plug the phone in and "
                 "go to sleep, that chip sits completely idle for eight hours.",
         "size": 11, "color": HAZE, "line": 1.32, "space_after": 9},
        {"text": "NeuroMesh rents out those idle hours. Phone owners earn credits overnight. "
                 "Developers get a big batch of AI work done by borrowing many phones at once, "
                 "instead of renting an expensive server.",
         "size": 11, "color": INK, "line": 1.32},
    ])

    cw = (CW - 0.5) / 3
    card(s, L, 3.95, cw, 2.45,
         label="Problem · the phone owner",
         heading="An expensive chip that earns nothing",
         body="Every flagship ships an NPU that sits dark from midnight to morning — charging, "
              "on Wi-Fi, screen off, cool. The most abundant idle AI chip on earth, and no way "
              "to make a single rupee from it.",
         plain="You already paid for the chip. It earns you nothing while you sleep.")
    card(s, L + cw + 0.25, 3.95, cw, 2.45,
         label="Problem · the developer",
         heading="Simple bulk work, priced like heavy work",
         body="A developer processing two million documents does not need a top-end AI server. "
              "They need throughput. They rent a scarce, costly GPU anyway, because nothing "
              "cheaper will take the job.",
         plain="Like hiring a crane to carry a thousand grocery bags.")
    card(s, L + 2 * (cw + 0.25), 3.95, cw, 2.45,
         label="The idea",
         heading="One app, two sides of a market",
         body="Provider mode lends the phone's NPU under strict conditions and earns credits. "
              "Consumer mode submits a batch, watches it split across the fleet, and pays in "
              "those same credits.",
         plain="Lend your phone at night, spend what it earned during the day.")

    x = L
    x = chip(s, x, 6.65, "✓  Working prototype already built", MINT, MINT_DIM)
    x = chip(s, x, 6.65, "6,474 lines of code")
    x = chip(s, x, 6.65, "30 automated tests, all passing")
    chip(s, x, 6.65, "Installable Android app, 43 MB")


def slide2(prs):
    s = slide_base(prs, 2, "What makes us different")

    title(s, "Networks like this already exist.\nEvery one of them rents computers. We rent phones.",
          y=0.82, size=23)

    table(s, L, 1.92, CW, 1.95, [
        ["NETWORK", "WHAT IT RENTS OUT", "WHO CAN SUPPLY IT", "WHAT YOU NEED TO JOIN"],
        ["Vast.ai", "Desktop & server graphics cards", "Crypto miners, small hosts",
         "A spare ₹1.5 lakh GPU"],
        ["Akash", "Data-centre graphics cards", "People who own server racks",
         "Infrastructure and uptime"],
        ["io.net", "Pooled GPU clusters", "Data centres, mining farms",
         "To already be a supplier"],
        ["Folding@home", "Home computer power", "Volunteers", "Nothing — but nobody is paid"],
        [("NeuroMesh",), ("Phone AI chips",), ("Anyone with a phone and a charger",),
         ("To install an app and go to bed",)],
    ], widths=[0.16, 0.27, 0.28, 0.29])

    cw = (CW - 0.5) / 3
    card(s, L, 4.12, cw, 2.30,
         label="Why now",
         heading='"Phones are too weak" stopped being true',
         body="That was correct when phone AI chips managed 5 trillion operations per second and "
              "overheated in ninety seconds. Today's flagships do 40–50 trillion. And the perfect "
              "conditions already happen in every home, every night, free.",
         plain="The technology caught up. Nobody updated the assumption.")
    card(s, L + cw + 0.25, 4.12, cw, 2.30,
         label="Our honest claim",
         heading="Not \"cheaper than a server\" — a closed loop",
         body="We do not pretend to beat data-centre pricing on raw cost. We offer what they "
              "cannot: earn credits while you sleep, spend them by day to run models your own "
              "phone is too small to run alone.",
         plain="The supplier and the customer are the same person.")
    card(s, L + 2 * (cw + 0.25), 4.12, cw, 2.30,
         label="Supply is the easy part",
         heading="The problem that kills these networks does not apply to us",
         body="Every compute marketplace struggles to find suppliers — they need people who own "
              "expensive hardware. Our supplier just needs a phone they already own.",
         plain="No new hardware has to exist anywhere for this to scale.")

    rect(s, L, 6.62, 0.022, 0.55, EMBER)
    textbox(s, L + 0.13, 6.62, CW - 0.2, 0.55, [{
        "text": "What we deliberately do not do, and why:  phones can run AI models but cannot "
                "train them — the chip is built for one and physically not the other. And we only "
                "accept work that splits into fully independent pieces, so no phone ever waits on "
                "another. Both limits are designed in on purpose.",
        "size": 8.5, "color": HAZE, "line": 1.25,
    }])


def slide3(prs):
    s = slide_base(prs, 3, "How it works")

    title(s, "One job, cut into slices, handed out one at a time "
             "so no two phones can take the same slice.", y=0.80, size=21)

    # glossary strip
    gy, gh = 1.58, 0.60
    box(s, L, gy, CW, gh, fill=PANEL)
    terms = [
        ("NPU", "The AI-only chip inside the phone."),
        ("SHARD", "One small slice of a big job — 8 items out of 120."),
        ("LEASE", "A 30-second reservation on a slice."),
        ("CREDITS", "The in-app currency: earned, then spent."),
    ]
    gw = CW / 4
    for i, (term, meaning) in enumerate(terms):
        gx = L + i * gw
        if i:
            rect(s, gx, gy + 0.08, 0.008, gh - 0.16, EDGE)
        textbox(s, gx + 0.16, gy + 0.10, gw - 0.30, 0.42, [
            {"text": term, "size": 8, "font": MONO, "bold": True, "color": MINT,
             "spacing": 1.1, "space_after": 2},
            {"text": meaning, "size": 8.5, "color": HAZE, "line": 1.15},
        ])

    # ---- diagram -------------------------------------------------------
    dy = 2.42

    # consumer phone
    box(s, 0.60, dy + 0.75, 1.45, 1.35)
    textbox(s, 0.72, dy + 0.92, 1.2, 1.0, [
        {"text": "Phone 1", "size": 11, "font": DISPLAY, "bold": True, "space_after": 3},
        {"text": "sends the job", "size": 8, "color": HAZE},
        {"text": "watches it live", "size": 8, "color": HAZE, "space_after": 5},
        {"text": '"consumer"', "size": 7.5, "font": MONO, "color": IRIS},
    ])

    # coordinator
    box(s, 2.35, dy, 5.05, 2.95, fill=None, line=EDGE)
    textbox(s, 2.50, dy + 0.10, 4.8, 0.2, [{
        "text": "THE COORDINATOR  (A SMALL SERVER)", "size": 7.5, "font": MONO,
        "color": HAZE, "spacing": 1.3}])

    box(s, 2.50, dy + 0.36, 4.75, 0.48, fill=PANEL2)
    textbox(s, 2.64, dy + 0.42, 4.5, 0.38, [
        {"text": "1 · Cut the job into slices", "size": 10, "font": DISPLAY, "bold": True,
         "space_after": 2},
        {"text": "price it · hold the credits · lock the data", "size": 8, "color": HAZE}])

    box(s, 2.50, dy + 0.92, 4.75, 0.85, fill=PANEL2)
    textbox(s, 2.64, dy + 0.97, 4.5, 0.2, [
        {"text": "2 · Queue by phone strength", "size": 10, "font": DISPLAY, "bold": True}])
    for j, (lbl, col) in enumerate([("strong phones · NPU", MINT_DIM),
                                    ("mid phones · graphics chip", IRIS_DIM),
                                    ("basic phones · main processor", IDLE)]):
        rect(s, 2.64, dy + 1.20 + j * 0.19, 4.45, 0.165, col)
        textbox(s, 2.74, dy + 1.222 + j * 0.19, 4.2, 0.14,
                [{"text": lbl, "size": 7.5, "font": MONO, "color": INK}])

    box(s, 2.50, dy + 1.85, 4.75, 0.50, fill=None, line=MINT, line_w=1.25)
    textbox(s, 2.64, dy + 1.90, 4.5, 0.4, [
        {"text": "3 · Hand out one slice", "size": 10, "font": DISPLAY, "bold": True,
         "color": MINT, "space_after": 2},
        {"text": "a single unbreakable step — no two phones can ever get the same slice",
         "size": 7.5, "color": HAZE}])

    box(s, 2.50, dy + 2.43, 2.30, 0.42, fill=PANEL2)
    textbox(s, 2.64, dy + 2.48, 2.1, 0.34, [
        {"text": "5 · Reassemble", "size": 9.5, "font": DISPLAY, "bold": True, "space_after": 1},
        {"text": "answers in the right order", "size": 7.5, "color": HAZE}])

    box(s, 4.95, dy + 2.43, 2.30, 0.42, fill=PANEL2)
    textbox(s, 5.09, dy + 2.48, 2.1, 0.34, [
        {"text": "6 · Pay out", "size": 9.5, "font": DISPLAY, "bold": True, "space_after": 1},
        {"text": "credits to each phone that helped", "size": 7.5, "color": HAZE}])

    # provider phones
    for j, (name, tier, tint, note) in enumerate([
        ("Phone 2", "NPU · strong", MINT_DIM, "4 · runs the slice"),
        ("Phone 3", "NPU · strong", MINT_DIM, "4 · runs the slice"),
        ("Phone 4", "graphics chip", IRIS_DIM, "backup path"),
    ]):
        py = dy + 0.10 + j * 1.00
        box(s, 8.55, py, 1.95, 0.85)
        textbox(s, 8.68, py + 0.09, 1.75, 0.7, [
            {"text": name, "size": 10, "font": DISPLAY, "bold": True, "space_after": 3}])
        rect(s, 8.68, py + 0.34, 0.95, 0.17, tint)
        textbox(s, 8.75, py + 0.36, 0.9, 0.14,
                [{"text": tier, "size": 7, "font": MONO, "color": INK}])
        textbox(s, 8.68, py + 0.58, 1.75, 0.16,
                [{"text": note, "size": 7.5, "color": HAZE}])

    # arrows
    line(s, 2.05, dy + 1.20, 2.48, dy + 0.62, HAZE, 1.0)
    textbox(s, 1.95, dy + 0.30, 1.3, 0.16,
            [{"text": '"here is my job"', "size": 7.5, "font": MONO, "color": HAZE}])

    line(s, 2.48, dy + 2.05, 2.07, dy + 1.72, IRIS, 1.0, dashed=True)
    textbox(s, 0.62, dy + 2.30, 1.6, 0.32, [
        {"text": "live progress feed", "size": 7.5, "font": MONO, "color": IRIS},
        {"text": "(no private data)", "size": 7.5, "font": MONO, "color": HAZE}])

    for j, ty in enumerate([dy + 0.50, dy + 1.50, dy + 2.50]):
        line(s, 7.42, dy + 2.05, 8.50, ty, MINT, 1.0)
    textbox(s, 7.48, dy + 1.05, 1.1, 0.30, [
        {"text": "one slice", "size": 7.5, "font": MONO, "color": MINT},
        {"text": "+ 30s hold", "size": 7.5, "font": MONO, "color": MINT}])

    line(s, 8.50, dy + 0.95, 7.30, dy + 2.60, HAZE, 0.9)
    textbox(s, 7.55, dy + 2.62, 1.5, 0.16,
            [{"text": "answer comes back", "size": 7.5, "color": HAZE}])

    line(s, 10.55, dy + 0.50, 10.55, dy + 2.45, EMBER, 1.0, dashed=True)
    textbox(s, 10.70, dy + 0.95, 2.0, 0.62, [
        {"text": "phone vanished?", "size": 8, "font": MONO, "color": EMBER},
        {"text": "the hold expires", "size": 8, "font": MONO, "color": EMBER, "space_after": 3},
        {"text": "→ slice goes back", "size": 8, "color": HAZE},
        {"text": "→ that phone earns 0", "size": 8, "color": HAZE}])

    # ---- four notes ----------------------------------------------------
    nw = (CW - 0.45) / 4
    notes = [
        ("Two kinds of memory",
         "A fast temporary store holds live work. A permanent database holds accounts and money.",
         "A whiteboard for today, a ledger book for the year."),
        ("Reservations, not trust",
         "Each slice is held for 30 seconds. A watchdog puts expired holds straight back in the queue.",
         "Quiet for 30 seconds? We assume the worst."),
        ("The money can't drift",
         "Whole numbers only. Exactly one piece of code may touch a balance, with an audit beside it.",
         "One cashier, one till, one audit."),
        ("Privacy & safety",
         "Data is locked until a phone holds the slice. Someone else's job returns \"not found\".",
         "\"Not allowed\" would confirm it exists."),
    ]
    for i, (lbl, body, plain) in enumerate(notes):
        card(s, L + i * (nw + 0.15), 5.62, nw, 1.55, label=lbl, body=body, plain=plain)


def slide4(prs):
    s = slide_base(prs, 4, "What already works")

    title(s, "We did not just design this. We built it, and then we tried to break it.",
          y=0.80, size=23)

    textbox(s, L, 1.42, 9.5, 0.24, [{
        "text": "Everything below is real output from the system running today — "
                "not a plan, not a mock-up.",
        "size": 10.5, "color": HAZE}])

    cw = (CW - 0.45) / 4
    stats = [
        ("A full job, start to finish", "15 / 15", MINT,
         "120 items cut into 15 slices, finished by 3 phones in 1.9 seconds. All 120 answers came "
         "back in the original order and the credit accounts balanced exactly."),
        ("We yanked a phone mid-job", "8 / 8", MINT,
         "A phone was pulled away holding slice #0. Its hold expired, the slice went back, another "
         "phone finished it. Nothing lost — and the phone that walked away was paid zero."),
        ("16 phones fought over 1 slice", "16 → 1", IRIS,
         "Exactly one won. Two phones on one slice means paying twice for one piece of work, so we "
         "made that impossible rather than unlikely."),
        ("Weak phones get protected", "0", IRIS,
         "Slices offered to an older phone from a job it could not handle. Each phone is judged on "
         "what it actually measures, never on its brand name."),
    ]
    for i, (lbl, stat, col, body) in enumerate(stats):
        card(s, L + i * (cw + 0.15), 1.82, cw, 2.05,
             label=lbl, stat=stat, stat_color=col, body=body)

    hw = (CW - 0.25) / 2
    card(s, L, 4.02, hw, 2.30,
         label="What the person who submitted the job sees")
    states = "dddddddccd" "dfdcdddd--"
    colors = {"d": MINT, "c": IRIS, "f": EMBER, "-": IDLE}
    for i, ch in enumerate(states):
        gx = L + 0.18 + (i % 10) * 0.27
        gy = 4.50 + (i // 10) * 0.27
        b = box(s, gx, gy, 0.22, 0.22, fill=colors[ch], line=None, radius=0.22)
    x = L + 0.18
    x = chip(s, x, 5.10, "finished", MINT, MINT_DIM)
    x = chip(s, x, 5.10, "running now", IRIS, IRIS_DIM)
    x = chip(s, x, 5.10, "dropped, retrying", EMBER, EMBER_DIM)
    chip(s, x, 5.10, "waiting")
    textbox(s, L + 0.18, 5.52, hw - 0.36, 0.6, [{
        "text": "One square per slice, updating live on the phone screen as the fleet works. "
                "This is the demo: the judges watch squares turn green in real time.",
        "size": 9, "color": HAZE, "line": 1.25}])

    card(s, L + hw + 0.25, 4.02, hw, 2.30,
         label="Built and verified today",
         bullets=[
             "The coordinator — ~20 web endpoints, live progress feed, automatic database setup.",
             "The scheduler — five scripts covering take, hold, finish, fail and recover.",
             "The Android app — one app, both modes, background service, real UI. 43 MB.",
             "Three real AI models bundled in, each fingerprint-checked before it may run.",
             "Test tooling — a fake fleet, including a switch that kills a phone mid-job.",
         ])

    rect(s, L, 6.50, 0.022, 0.62, EMBER)
    textbox(s, L + 0.13, 6.50, CW - 0.2, 0.62, [{
        "text": "What we have NOT proven yet, said plainly:  the Android app compiles and installs, "
                "but we have never run it on an actual iQOO 15 — we do not own one. Getting real "
                "inference running on the real chip is hour 0–3 of the hackathon, and it is exactly "
                "why we want to be in the room.",
        "size": 9, "color": HAZE, "line": 1.25}])


def slide5(prs):
    s = slide_base(prs, 5, "Phone-first · On-device AI · Performance")

    title(s, "The phone is not where we show the project. The phone is the project.",
          y=0.80, size=23)

    cw = (CW - 0.5) / 3
    card(s, L, 1.45, cw, 3.05,
         label="Which chip runs the model",
         heading="Try the AI chip, then the graphics chip, then the processor",
         body="A phone has three chips that can run a model, best to worst. We ask for Qualcomm's "
              "AI chip first, fall back to the graphics chip, then the main processor. If one "
              "fails to start we quietly step down instead of crashing.",
         plain="A slow demo beats a dead demo. Something always runs.")
    card(s, L + cw + 0.25, 1.45, cw, 3.05,
         label="The models we ship",
         heading="Three real AI models, shrunk to fit a phone",
         bullets=[
             "Text understanding — turns a sentence into comparable numbers. 188 KB.",
             "Image recognition — real trained classifier, 1,000 categories. 3.9 MB.",
             "Settings search — scores AI configurations to find the best one.",
         ],
         plain="We tested two versions of the photo model and kept the one that did not break "
               "on older phones.")
    card(s, L + 2 * (cw + 0.25), 1.45, cw, 3.05,
         label="Respecting the owner",
         heading="The phone only works when it can genuinely spare the effort",
         body="Charging AND on free Wi-Fi AND screen off AND not warm. All four, re-checked before "
              "every single slice, on the phone and again on the server. We stop before the phone "
              "gets hot enough to notice — not after.",
         plain="Pick up your phone and it stops instantly. Your battery, data and evening are "
               "never touched.")

    hw = (CW - 0.25) / 2
    card(s, L, 4.72, hw, 1.75,
         label="Why this is genuinely phone-first",
         bullets=[
             "The phone is the engine — take the phones away and there is no product.",
             "The phone is even the server. No laptop: a script runs the whole coordinator on one iQOO 15.",
             "The phone is the control panel — submitting, watching, earnings, all on-device.",
             "The phone set the rules. The 30-second hold and heat limit were chosen for phones.",
         ])
    card(s, L + hw + 0.25, 4.72, hw, 1.75,
         label="What we will measure at the hackathon",
         bullets=[
             "Real numbers from the real chip — same slice on AI chip vs graphics chip vs processor, "
             "with speed and battery cost published.",
             "An hour-long heat test to tune our temperature limit against how the iQOO 15 really behaves.",
             "A big-screen judging view: mirror the live progress grid using vivo's multi-screen feature.",
         ])


def slide6(prs):
    s = slide_base(prs, 6, "The 30 hours · Growth · Team")

    title(s, "Most teams will spend the first day building the hard middle. Ours is already done.",
          y=0.80, size=23)

    table(s, L, 1.50, CW, 2.45, [
        ["HOURS", "THE MILESTONE", "WHERE WE ALREADY STAND", "SO WE SPEND THE TIME ON"],
        ["0–3", "Get an AI model running on the phone itself", "Not tried on real hardware",
         "First real run on the iQOO 15. Our only unproven step."],
        ["3–8", "Background service that fetches and runs work", "Written, not device-tested",
         "Prove the backup path on hardware, then the AI chip."],
        [("8–14",), ("Splitting, scheduling, reassembling, paying",), ("BUILT AND TESTED",),
         ("Six hours we get back — measuring the chip, tuning for heat.",)],
        ["14–20", "End to end: one phone, then many", "Proven with fake phones",
         "Swap the fake fleet for four real ones. Nothing else changes."],
        ["20–26", "Job-submission screen and live progress", "Built",
         "Polish, the big-screen judging view, earnings history."],
        ["26–30", "Rehearse, break it on purpose, keep a buffer", "Break-test automated",
         "Yank a phone on stage, on purpose, and let it recover live."],
    ], widths=[0.07, 0.28, 0.22, 0.43])

    cw = (CW - 0.5) / 3
    card(s, L, 4.15, cw, 2.35,
         label="Can it grow?",
         heading="Designed for a crowd, not for four phones",
         body="Handing out a slice costs the same whether ten phones are asking or ten thousand. "
              "The server never tracks who is alive — the hold clock does it free. Going from 4 "
              "phones to 4,000 means a bigger queue and more copies of the same server.",
         plain="Nothing in the design gets slower as more people join.")
    card(s, L + cw + 0.25, 4.15, cw, 2.35,
         label="Who would actually use it",
         heading="Both sides are paid in the same currency",
         body="A student processes a research dataset overnight using credits their own phone "
              "earned while charging. A phone owner turns eight dead hours into something "
              "spendable. Where renting a GPU costs real money and a phone is already in every "
              "pocket, that swap matters.")
    card(s, L + 2 * (cw + 0.25), 4.15, cw, 2.35,
         label="Team  ← replace these",
         heading="Who is building it",
         bullets=[
             "[Name] — Android & on-device AI: the app, the chips, the background service.",
             "[Name] — Server & distributed systems: scheduling, holds, credits.",
             "[Name] — Product & interface: the screens, the live view, the demo.",
             "[Name] — Infrastructure & testing: databases, break-tests, rehearsal.",
         ])

    rect(s, L, 6.70, CW, 0.014, EDGE)
    textbox(s, L, 6.92, CW, 0.3, [{
        "text": "NeuroMesh · renting out the AI chip in your pocket        "
                "GitHub: [repo link]        Demo video: [link]        "
                "Full technical write-up: README.md",
        "size": 8.5, "font": MONO, "color": HAZE}])


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    for build in (slide1, slide2, slide3, slide4, slide5, slide6):
        build(prs)

    prs.save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
